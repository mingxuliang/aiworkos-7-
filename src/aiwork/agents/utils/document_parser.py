# -*- coding: utf-8 -*-
"""Document text extraction for common file types.

Supports: PDF, DOCX, XLSX/XLS, PPTX, TXT, CSV, MD, and image OCR fallback.
All parsers are imported lazily so missing optional deps degrade gracefully.
"""
from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

# Max characters extracted per document to avoid flooding the LLM context.
MAX_EXTRACT_CHARS = 80_000

# Extensions handled by each parser
_PDF_EXTS = {".pdf"}
_WORD_EXTS = {".docx", ".doc"}
_EXCEL_EXTS = {".xlsx", ".xls", ".xlsm"}
_PPT_EXTS = {".pptx", ".ppt"}
_TEXT_EXTS = {".txt", ".md", ".rst", ".csv", ".tsv", ".json", ".yaml", ".yml",
              ".xml", ".html", ".htm", ".log", ".py", ".js", ".ts", ".java",
              ".c", ".cpp", ".h", ".go", ".rs", ".sh", ".sql", ".toml", ".ini",
              ".cfg", ".conf"}
_IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tiff", ".tif"}


def _ext(path: str) -> str:
    return Path(path).suffix.lower()


# ---------------------------------------------------------------------------
# Individual parsers
# ---------------------------------------------------------------------------

def _parse_pdf(path: str) -> Optional[str]:
    try:
        import fitz  # PyMuPDF
    except ImportError:
        logger.warning("PyMuPDF not installed; PDF text extraction unavailable")
        return None

    try:
        doc = fitz.open(path)
        parts: list[str] = []
        for page in doc:
            text = page.get_text("text")
            if text.strip():
                parts.append(text)
        doc.close()
        return "\n".join(parts) or None
    except Exception as e:
        logger.warning("PDF parse error for %s: %s", path, e)
        return None


def _parse_docx(path: str) -> Optional[str]:
    try:
        from docx import Document  # python-docx
    except ImportError:
        return None

    try:
        doc = Document(path)
        lines: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                lines.append(para.text)
        # Also extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    lines.append(row_text)
        return "\n".join(lines) or None
    except Exception as e:
        logger.warning("DOCX parse error for %s: %s", path, e)
        return None


def _parse_xlsx(path: str) -> Optional[str]:
    try:
        import openpyxl
    except ImportError:
        return None

    try:
        wb = openpyxl.load_workbook(path, data_only=True, read_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(c.strip() for c in cells):
                    parts.append("\t".join(cells))
        wb.close()
        return "\n".join(parts) or None
    except Exception as e:
        logger.warning("XLSX parse error for %s: %s", path, e)
        return None


def _parse_xls(path: str) -> Optional[str]:
    try:
        import xlrd  # type: ignore
    except ImportError:
        return None

    try:
        wb = xlrd.open_workbook(path)
        parts: list[str] = []
        for sheet in wb.sheets():
            parts.append(f"=== Sheet: {sheet.name} ===")
            for r in range(sheet.nrows):
                cells = [str(sheet.cell_value(r, c)) for c in range(sheet.ncols)]
                if any(c.strip() for c in cells):
                    parts.append("\t".join(cells))
        return "\n".join(parts) or None
    except Exception as e:
        logger.warning("XLS parse error for %s: %s", path, e)
        return None


def _parse_pptx(path: str) -> Optional[str]:
    try:
        from pptx import Presentation  # python-pptx
    except ImportError:
        return None

    try:
        prs = Presentation(path)
        slides: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                # Tables inside shapes
                if shape.has_table:
                    for row in shape.table.rows:
                        row_text = "\t".join(
                            cell.text.strip() for cell in row.cells
                        )
                        if row_text.strip():
                            texts.append(row_text)
            if texts:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
        return "\n\n".join(slides) or None
    except Exception as e:
        logger.warning("PPTX parse error for %s: %s", path, e)
        return None


def _parse_text(path: str) -> Optional[str]:
    encodings = ["utf-8", "utf-8-sig", "gbk", "gb18030", "latin-1"]
    for enc in encodings:
        try:
            text = Path(path).read_text(encoding=enc)
            return text or None
        except (UnicodeDecodeError, LookupError):
            continue
        except Exception as e:
            logger.warning("Text read error for %s: %s", path, e)
            return None
    return None


def _parse_image_ocr(path: str) -> Optional[str]:
    """Attempt OCR on an image using RapidOCR (no system deps required)."""
    try:
        from rapidocr_onnxruntime import RapidOCR  # type: ignore
        ocr = RapidOCR()
        result, _ = ocr(path)
        if result:
            lines = [item[1] for item in result if item and len(item) > 1]
            return "\n".join(lines) or None
        return None
    except ImportError:
        pass
    except Exception as e:
        logger.warning("RapidOCR error for %s: %s", path, e)

    # Fallback: pytesseract if available
    try:
        import pytesseract  # type: ignore
        from PIL import Image
        img = Image.open(path)
        text = pytesseract.image_to_string(img, lang="chi_sim+eng")
        return text.strip() or None
    except ImportError:
        pass
    except Exception as e:
        logger.warning("pytesseract error for %s: %s", path, e)

    return None


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def extract_text_from_file(path: str) -> Optional[str]:
    """Extract plain text from a file based on its extension.

    Returns the extracted text (truncated to MAX_EXTRACT_CHARS) or None if
    the file type is unsupported or extraction fails.
    """
    if not os.path.isfile(path):
        return None

    ext = _ext(path)

    if ext in _PDF_EXTS:
        text = _parse_pdf(path)
    elif ext in _WORD_EXTS:
        if ext == ".docx":
            text = _parse_docx(path)
        else:
            # .doc — try docx fallback (may fail for old binary format)
            text = _parse_docx(path)
    elif ext in _EXCEL_EXTS:
        if ext == ".xls":
            text = _parse_xls(path)
        else:
            text = _parse_xlsx(path)
    elif ext in _PPT_EXTS:
        text = _parse_pptx(path)
    elif ext in _TEXT_EXTS:
        text = _parse_text(path)
    elif ext in _IMAGE_EXTS:
        text = _parse_image_ocr(path)
    else:
        return None

    if not text:
        return None

    if len(text) > MAX_EXTRACT_CHARS:
        text = text[:MAX_EXTRACT_CHARS] + f"\n\n[... content truncated at {MAX_EXTRACT_CHARS} chars ...]"

    return text


def is_parseable(path: str) -> bool:
    """Return True if extract_text_from_file can handle this file extension."""
    ext = _ext(path)
    return ext in (
        _PDF_EXTS | _WORD_EXTS | _EXCEL_EXTS | _PPT_EXTS | _TEXT_EXTS | _IMAGE_EXTS
    )
